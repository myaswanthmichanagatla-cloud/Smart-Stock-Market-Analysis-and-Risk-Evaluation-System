#!/usr/bin/env python3
"""
Stock Market RAG Presentation Generator
Creates a highly polished, creative presentation matching the custom light geometric network template.
"""

import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

# Reconfigure standard streams to prevent Windows CP1252 encoding crashes
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8')

# Design Tokens (Matching the User's Light Geometric Network Template)
BG_COLOR = RGBColor(248, 249, 250)       # Clean Off-white Canvas
TEXT_DARK = RGBColor(26, 36, 43)         # Deep Charcoal/Black
TEXT_MUTED = RGBColor(100, 116, 139)     # Soft Slate Gray
CYAN_ACCENT = RGBColor(34, 186, 250)     # Sky Blue / Cyan Accent
NAVY_ACCENT = RGBColor(23, 42, 69)       # Dark Navy
CARD_BG = RGBColor(238, 242, 246)        # Card Container Fill
DARK_BOX_BG = RGBColor(34, 38, 41)       # Cover slide block / Terminal Background
WHITE = RGBColor(255, 255, 255)
RED_ACCENT = RGBColor(239, 68, 68)       # Warm Red for Problems
GREEN_ACCENT = RGBColor(34, 197, 94)     # Success Green

# -------------------------------------------------------------
# BACKGROUND DECORATION: GEOMETRIC NEURAL NETWORK MESH
# -------------------------------------------------------------
def draw_network_background(slide, slide_width, slide_height, mesh_type="full"):
    """
    Programmatically draws the geometric network mesh background using thin lines and circles
    to match the user's template screenshot exactly.
    """
    if mesh_type == "full":
        # Nodes: (x, y, radius, type, color)
        nodes = [
            (Inches(0.8), Inches(1.2), Inches(0.12), "solid", TEXT_DARK),
            (Inches(1.8), Inches(0.8), Inches(0.1), "solid", CYAN_ACCENT),
            (Inches(3.2), Inches(1.5), Inches(0.08), "solid", TEXT_MUTED),
            (Inches(5.5), Inches(0.9), Inches(0.24), "nested", TEXT_DARK), 
            (Inches(6.2), Inches(1.7), Inches(0.1), "solid", CYAN_ACCENT),
            (Inches(7.5), Inches(1.1), Inches(0.15), "solid", TEXT_MUTED),
            (Inches(2.2), Inches(3.2), Inches(0.3), "nested", CYAN_ACCENT),
            (Inches(3.8), Inches(2.5), Inches(0.14), "solid", TEXT_DARK),
            (Inches(5.1), Inches(2.8), Inches(0.12), "solid", TEXT_MUTED),
            (Inches(6.8), Inches(2.4), Inches(0.18), "solid", CARD_BG),
            (Inches(8.0), Inches(2.9), Inches(0.1), "solid", TEXT_DARK),
            (Inches(1.1), Inches(4.8), Inches(0.1), "solid", TEXT_DARK),
            (Inches(3.0), Inches(4.5), Inches(0.12), "solid", CYAN_ACCENT),
            (Inches(4.5), Inches(4.0), Inches(0.08), "solid", TEXT_MUTED),
            (Inches(5.8), Inches(4.6), Inches(0.15), "solid", CARD_BG),
            (Inches(7.2), Inches(4.2), Inches(0.12), "solid", CYAN_ACCENT),
            (Inches(8.8), Inches(4.5), Inches(0.22), "nested", TEXT_DARK),
            (Inches(10.2), Inches(1.5), Inches(0.12), "solid", TEXT_MUTED),
            (Inches(11.5), Inches(0.8), Inches(0.15), "solid", CYAN_ACCENT),
            (Inches(12.2), Inches(2.2), Inches(0.25), "nested", TEXT_DARK),
            (Inches(10.8), Inches(3.2), Inches(0.1), "solid", CYAN_ACCENT),
            (Inches(12.5), Inches(4.2), Inches(0.12), "solid", TEXT_MUTED),
        ]
        connections = [
            (0, 1), (0, 6), (1, 2), (1, 3), (2, 3), (2, 7), (3, 4), (4, 5), (4, 9), (5, 10),
            (6, 7), (6, 11), (7, 8), (7, 12), (8, 9), (8, 13), (9, 10), (9, 14), (11, 12),
            (12, 13), (13, 14), (14, 15), (15, 16), (10, 16), (16, 20), (17, 18), (17, 19),
            (18, 19), (19, 20), (20, 21)
        ]
    elif mesh_type == "corner":
        # Subtly placed corner network mesh for content slides
        nodes = [
            (Inches(11.0), Inches(0.6), Inches(0.08), "solid", TEXT_MUTED),
            (Inches(11.6), Inches(0.4), Inches(0.06), "solid", CYAN_ACCENT),
            (Inches(12.3), Inches(0.7), Inches(0.15), "nested", TEXT_DARK),
            (Inches(11.8), Inches(1.2), Inches(0.08), "solid", TEXT_DARK),
            (Inches(12.6), Inches(1.4), Inches(0.06), "solid", CYAN_ACCENT),
        ]
        connections = [
            (0, 1), (1, 2), (0, 3), (2, 3), (2, 4), (3, 4)
        ]
    else:
        return

    # Draw lines (underneath nodes)
    for start, end in connections:
        n1, n2 = nodes[start], nodes[end]
        x1 = n1[0] + n1[2] / 2
        y1 = n1[1] + n1[2] / 2
        x2 = n2[0] + n2[2] / 2
        y2 = n2[1] + n2[2] / 2
        
        try:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
            conn.line.color.rgb = RGBColor(210, 216, 222)
            conn.line.width = Pt(0.75)
        except Exception:
            pass

    # Draw circle shapes
    for x, y, size, s_type, color in nodes:
        if s_type == "solid":
            c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
            c.fill.solid()
            c.fill.fore_color.rgb = color
            c.line.fill.background()
        elif s_type == "nested":
            # Outer hollow ring
            out_sz = size * 1.5
            offset = (out_sz - size) / 2
            out = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - offset, y - offset, out_sz, out_sz)
            out.fill.background()
            out.line.color.rgb = RGBColor(180, 186, 192)
            out.line.width = Pt(1)
            
            # Inner solid circle
            inn = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
            inn.fill.solid()
            inn.fill.fore_color.rgb = color
            inn.line.fill.background()

# -------------------------------------------------------------
# HELPER: CREATE BASE SLIDE WITH BG & TITLE
# -------------------------------------------------------------
def create_base_slide(prs, title_text):
    """Creates a base slide with a clean canvas background and elegant header."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Fill solid background color
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Draw top-right corner network graphic
    draw_network_background(slide, Inches(13.333), Inches(7.5), "corner")
    
    # Elegant Title text
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.0), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.name = "Arial Black"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    # Small Cyan Accent Line under title
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.1), Inches(1.5), Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN_ACCENT
    accent.line.fill.background()
    
    return slide

# -------------------------------------------------------------
# HELPER: ADD CARD / BOX CONTAINER
# -------------------------------------------------------------
def add_card(slide, left, top, width, height, title, content_list, accent_color=CYAN_ACCENT):
    """Adds a structured information card to a slide."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(218, 224, 230)
    card.line.width = Pt(1.5)
    
    # Text Frame
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.15)
    tf.margin_left = tf.margin_right = Inches(0.2)
    
    p_title = tf.paragraphs[0]
    p_title.text = title.upper()
    p_title.font.name = "Arial Black"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK
    p_title.space_after = Pt(8)
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)
        
    return card

# -------------------------------------------------------------
# SLIDE 1: TITLE & COVER SLIDE
# -------------------------------------------------------------
def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Solid off-white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Draw geometric neural network cover background
    draw_network_background(slide, Inches(13.333), Inches(7.5), "full")
    
    # Dark Charcoal Title Block (Bottom-Right)
    title_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(7.2), Inches(2.8))
    title_block.fill.solid()
    title_block.fill.fore_color.rgb = DARK_BOX_BG
    title_block.line.fill.background()
    
    # Inner thin white border matching user's template screenshot
    inner_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(4.0), Inches(6.8), Inches(2.4))
    inner_block.fill.background()
    inner_block.line.color.rgb = WHITE
    inner_block.line.width = Pt(1)
    
    # Title Text Frame
    tf_box = slide.shapes.add_textbox(Inches(5.8), Inches(4.1), Inches(6.6), Inches(2.2))
    tf = tf_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    
    p_title = tf.paragraphs[0]
    p_title.text = "SMART STOCK MARKET ANALYSIS\n& RISK EVALUATION SYSTEM"
    p_title.font.name = "Arial Black"
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.space_after = Pt(10)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "AI-POWERED SEMANTIC RAG & MACHINE LEARNING SYSTEM"
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(11)
    p_sub.font.bold = True
    p_sub.font.color.rgb = CYAN_ACCENT
    p_sub.space_after = Pt(12)
    
    p_team = tf.add_paragraph()
    p_team.text = "PRESENTED BY TEAM MEMBERS:\nSE25MAID035 | SE25MAID014 | SE25MAID016 | SE25MAID021"
    p_team.font.name = "Calibri"
    p_team.font.size = Pt(10)
    p_team.font.color.rgb = RGBColor(200, 200, 200)

# -------------------------------------------------------------
# SLIDE 2: THE CORE PROBLEM
# -------------------------------------------------------------
def build_slide_2(prs):
    slide = create_base_slide(prs, "The Financial Analysis Problem")
    
    # Left Card: Legacy Stock Screeners (Red Alert Accent)
    card_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = RED_ACCENT
    card_left.line.width = Pt(2)
    
    tf_l = card_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = Inches(0.3)
    tf_l.margin_top = Inches(0.3)
    
    p = tf_l.paragraphs[0]
    p.text = "THE LEGACY LIMITATIONS"
    p.font.name = "Arial Black"
    p.font.size = Pt(15)
    p.font.color.rgb = RED_ACCENT
    p.space_after = Pt(15)
    
    bullet_l = [
        "Keyword Restrictions: Legacy screeners rely strictly on database text fields, missing semantic intent.",
        "Query Failure: Searching for 'high yield banking stocks with negative sentiment buffer' fails on traditional filters.",
        "Isolated Data silos: Technical analysis, sentiment scoring, and portfolio risk are analyzed independently.",
        "Predictive Gaps: Lack of real-time machine learning prediction models integrated directly within financial indexing."
    ]
    for b in bullet_l:
        p = tf_l.add_paragraph()
        p.text = "• " + b
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)
        
    # Right Card: Modern AI RAG Solution (Green Accent)
    card_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = CYAN_ACCENT
    card_right.line.width = Pt(2)
    
    tf_r = card_right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.3)
    tf_r.margin_top = Inches(0.3)
    
    p = tf_r.paragraphs[0]
    p.text = "THE StockRAG INNOVATION"
    p.font.name = "Arial Black"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(15)
    
    bullet_r = [
        "Semantic Vector Search: FinBERT/SentenceTransformers map query concepts directly to financial metadata.",
        "Natural Language Interface: Translates complex user investment goals into precise vector queries.",
        "Contextual Risk Retrieval: Combines real-time volatility calculations, expected returns, and sentiment indexes.",
        "Integrated Predictors: Combines baseline linear predictions and XGBoost analysis directly inside retrieval loops."
    ]
    for b in bullet_r:
        p = tf_r.add_paragraph()
        p.text = "• " + b
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)

# -------------------------------------------------------------
# SLIDE 3: PROJECT OBJECTIVES
# -------------------------------------------------------------
def build_slide_3(prs):
    slide = create_base_slide(prs, "Project Scope & Objectives")
    
    add_card(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.4), 
             "1. Semantic RAG", 
             [
                 "Design vector-based indexing of stock news and data.",
                 "Implement FinBERT/Transformers semantic embeddings.",
                 "Calculate accurate Cosine Similarity rankings.",
                 "Bridge conversational finance to raw database fields."
             ])
             
    add_card(slide, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.4), 
             "2. Risk Evaluation", 
             [
                 "Automated extraction of financial performance KPIs.",
                 "Calculate daily return, annualized volatility, and beta.",
                 "Generate risk indicators (Low, Moderate, High risk).",
                 "Perform Modern Portfolio Theory optimal allocations."
             ], accent_color=CYAN_ACCENT)
             
    add_card(slide, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.4), 
             "3. Interactive UI", 
             [
                 "Establish Python-backed Twelve Data proxy server.",
                 "Build modern visual dashboard in pure HTML/JS.",
                 "Dynamic charting for stock returns and allocations.",
                 "Interactive RAG semantic query search console."
             ])

# -------------------------------------------------------------
# SLIDE 4: SYSTEM ARCHITECTURE FLOWCHART
# -------------------------------------------------------------
def build_slide_4(prs):
    slide = create_base_slide(prs, "System Architecture")
    
    # 4 Flowchart Process Cards horizontally connected
    steps = [
        ("DATA PIPELINE", ["yfinance API Loader", "Feature Engineering", "CSV Preprocessing"], Inches(0.8)),
        ("VECTOR & EMBEDDING", ["FinBERT Encoding", "Cosine Similarity Engine", "Pickled Stock Index"], Inches(3.8)),
        ("PREDICTION SYSTEM", ["Baseline LR Model", "XGBoost ML Model", "Portfolio Optimization"], Inches(6.8)),
        ("USER DASHBOARD", ["Flask API Proxy Server", "Interactive Charting", "HTML/JS Client UI"], Inches(9.8)),
    ]
    
    # Draw horizontal blocks
    for title, list_items, left in steps:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(2.6), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CYAN_ACCENT
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = tf.margin_right = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial Black"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
        
        for item in list_items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.name = "Calibri"
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(8)
            
    # Draw horizontal connection arrows between steps
    arrows = [
        (Inches(3.4), Inches(4.25), Inches(3.8)),
        (Inches(6.4), Inches(4.25), Inches(6.8)),
        (Inches(9.4), Inches(4.25), Inches(9.8))
    ]
    for start_x, y, end_x in arrows:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, y, end_x, y)
        conn.line.color.rgb = CYAN_ACCENT
        conn.line.width = Pt(2.5)

# -------------------------------------------------------------
# SLIDE 5: DETAILED RAG PIPELINE FLOWCHART
# -------------------------------------------------------------
def build_slide_5(prs):
    slide = create_base_slide(prs, "Detailed Financial RAG Pipeline")
    
    # Let's draw a complete flowchart with boxes, diamonds and connectors
    def add_flow_box(text, left, top, w, h, s_type=MSO_SHAPE.ROUNDED_RECTANGLE, bg=CARD_BG, border=CYAN_ACCENT):
        s = slide.shapes.add_shape(s_type, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = bg
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        return s

    # Nodes
    add_flow_box("USER QUERY\n(Natural Language)", Inches(0.6), Inches(3.2), Inches(1.8), Inches(1.0), MSO_SHAPE.OVAL, DARK_BOX_BG, CYAN_ACCENT)
    # text colors
    slide.shapes[-1].text_frame.paragraphs[0].font.color.rgb = WHITE
    
    add_flow_box("FINBERT\nTRANSFORMER\n(Generate Embedding)", Inches(2.8), Inches(3.2), Inches(1.8), Inches(1.0))
    add_flow_box("COSINE SIMILARITY\nSEARCH\n(vs Vector Store)", Inches(5.0), Inches(3.2), Inches(1.8), Inches(1.0))
    
    # Decision Diamond
    add_flow_box("MATCH FOUND?", Inches(7.2), Inches(2.9), Inches(1.6), Inches(1.6), MSO_SHAPE.DIAMOND, CARD_BG, RGBColor(245, 158, 11))
    
    add_flow_box("RETRIEVE TOP-K\nMETADATA\n(PE, Sentiment, Beta)", Inches(9.2), Inches(3.2), Inches(1.8), Inches(1.0))
    add_flow_box("DASHBOARD SYSTEM\nRESPONSE\n(Synthesized Report)", Inches(11.3), Inches(3.2), Inches(1.5), Inches(1.0), MSO_SHAPE.OVAL, DARK_BOX_BG, CYAN_ACCENT)
    slide.shapes[-1].text_frame.paragraphs[0].font.color.rgb = WHITE

    # Secondary Path (Fallback)
    add_flow_box("DEFAULT SECTOR\nBACKUP SCAN", Inches(7.1), Inches(5.2), Inches(1.8), Inches(0.8), MSO_SHAPE.ROUNDED_RECTANGLE, CARD_BG, RED_ACCENT)

    # Connections
    def arrow(x1, y1, x2, y2, color=CYAN_ACCENT):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = color
        conn.line.width = Pt(1.5)

    arrow(Inches(2.4), Inches(3.7), Inches(2.8), Inches(3.7))
    arrow(Inches(4.6), Inches(3.7), Inches(5.0), Inches(3.7))
    arrow(Inches(6.8), Inches(3.7), Inches(7.2), Inches(3.7))
    arrow(Inches(8.8), Inches(3.7), Inches(9.2), Inches(3.7))
    arrow(Inches(11.0), Inches(3.7), Inches(11.3), Inches(3.7))
    
    # Decision paths
    arrow(Inches(8.0), Inches(4.5), Inches(8.0), Inches(5.2), RGBColor(245, 158, 11))
    arrow(Inches(8.9), Inches(5.6), Inches(10.1), Inches(5.6))
    arrow(Inches(10.1), Inches(5.6), Inches(10.1), Inches(4.2))
    
    # Yes / No Text
    txt_yes = slide.shapes.add_textbox(Inches(8.8), Inches(3.4), Inches(0.4), Inches(0.3))
    txt_yes.text_frame.paragraphs[0].text = "Yes"
    txt_yes.text_frame.paragraphs[0].font.size = Pt(9)
    txt_yes.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

    txt_no = slide.shapes.add_textbox(Inches(8.1), Inches(4.6), Inches(0.4), Inches(0.3))
    txt_no.text_frame.paragraphs[0].text = "No"
    txt_no.text_frame.paragraphs[0].font.size = Pt(9)
    txt_no.text_frame.paragraphs[0].font.color.rgb = RED_ACCENT

# -------------------------------------------------------------
# SLIDE 6: DATA PIPELINE & PREPROCESSING FLOWCHART
# -------------------------------------------------------------
def build_slide_6(prs):
    slide = create_base_slide(prs, "Data Preprocessing Pipeline")
    
    # Let's draw a vertical/grid layout flowchart for yfinance pipeline
    def add_node(text, left, top, w, h, border=CYAN_ACCENT):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = CARD_BG
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        return s

    add_node("RAW DATA INGESTION\nyfinance Stock Data Extract", Inches(1.5), Inches(2.0), Inches(3.2), Inches(1.0))
    add_node("DATA CLEANING\nForward fill, NaN pruning", Inches(1.5), Inches(3.6), Inches(3.2), Inches(1.0))
    add_node("FEATURE ENGINEERING\n40+ technical metrics (RSI, EMA)", Inches(1.5), Inches(5.2), Inches(3.2), Inches(1.0))
    
    add_node("SENTIMENT ANALYZER\nNews classification & scores", Inches(6.5), Inches(2.0), Inches(3.2), Inches(1.0))
    add_node("RISK COMPLIANCE METADATA\nVolatility, Beta, Max Drawdowns", Inches(6.5), Inches(3.6), Inches(3.2), Inches(1.0))
    add_node("RAG INDEX CREATION\npickle metadata (.json + .pkl)", Inches(6.5), Inches(5.2), Inches(3.2), Inches(1.0))

    # Draw connection lines between nodes
    def arrow(x1, y1, x2, y2):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = CYAN_ACCENT
        conn.line.width = Pt(2)

    arrow(Inches(3.1), Inches(3.0), Inches(3.1), Inches(3.6))
    arrow(Inches(3.1), Inches(4.6), Inches(3.1), Inches(5.2))
    
    # Bridge to second column
    arrow(Inches(4.7), Inches(5.7), Inches(6.5), Inches(5.7))
    
    arrow(Inches(8.1), Inches(5.2), Inches(8.1), Inches(4.6))
    arrow(Inches(8.1), Inches(3.6), Inches(8.1), Inches(3.0))

# -------------------------------------------------------------
# SLIDE 7: TECHNICAL STACK
# -------------------------------------------------------------
def build_slide_7(prs):
    slide = create_base_slide(prs, "System Technology Stack")
    
    # 4 Quadrants
    add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.3),
             "MODELS & EMBEDDINGS",
             [
                 "FinBERT (Financial sentiment processing)",
                 "SentenceTransformers (All-MiniLM-L6-v2)",
                 "scikit-learn (Baseline Predictor)",
                 "XGBoost Classifier (Advanced Model)"
             ])

    add_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.3),
             "CORE STACK & DATA",
             [
                 "Python 3.12 (Core pipeline)",
                 "yfinance (Automatic historical market ingestion)",
                 "Pandas & NumPy (Feature extraction, data structures)",
                 "Scipy (Advanced statistics & VaR metrics)"
             ], accent_color=CYAN_ACCENT)

    add_card(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.3),
             "WEB DASHBOARD & PROXY",
             [
                 "Python Flask proxy server (Local API wrapper)",
                 "HTML5 & Vanilla CSS (Premium sleek UI layout)",
                 "JavaScript (Dynamic updates and async fetches)",
                 "Twelve Data API integration (External quotes)"
             ], accent_color=CYAN_ACCENT)

    add_card(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.3),
             "INDEX & DEPLOYMENT",
             [
                 "Pickle Indexing (RAG vector knowledge base)",
                 "JSON formats (Inter-process dataset exchanges)",
                 "Localhost server execution environment",
                 "PowerShell/Windows shell automated runners"
             ])

# -------------------------------------------------------------
# SLIDE 8: DATASET & SEMANTIC INDEX OVERVIEW
# -------------------------------------------------------------
def build_slide_8(prs):
    slide = create_base_slide(prs, "Dataset & RAG Index Metrics")
    
    # Big statistics KPI text
    def add_kpi(slide, left, top, num, label, sub):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.6), Inches(4.4))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = RGBColor(218, 224, 230)
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.4)
        
        p_num = tf.paragraphs[0]
        p_num.text = num
        p_num.alignment = PP_ALIGN.CENTER
        p_num.font.name = "Arial Black"
        p_num.font.size = Pt(36)
        p_num.font.bold = True
        p_num.font.color.rgb = CYAN_ACCENT
        p_num.space_after = Pt(15)
        
        p_lbl = tf.add_paragraph()
        p_lbl.text = label.upper()
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.font.name = "Arial Black"
        p_lbl.font.size = Pt(12)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = TEXT_DARK
        p_lbl.space_after = Pt(15)
        
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = TEXT_MUTED

    add_kpi(slide, Inches(0.8), Inches(1.8), "66,000+", "Total Data records", 
            "Historical stock prices, closing prices, and trading volume calculated across major market indexes.")
            
    add_kpi(slide, Inches(4.8), Inches(1.8), "12+", "Sectors Analyzed", 
            "Broad sector diversification including Technology, Banking, Healthcare, Energy, and Financial Indexes.")
            
    add_kpi(slide, Inches(8.8), Inches(1.8), "5,000", "Semantic Index", 
            "Knowledge base vector representations generated using SentenceTransformers to map sentiment, risk and expected return.")

# -------------------------------------------------------------
# SLIDE 9: RAG DEMO TERMINAL MOCKUP
# -------------------------------------------------------------
def build_slide_9(prs):
    slide = create_base_slide(prs, "Semantic RAG Query in Action")
    
    # Terminal block background
    term = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    term.fill.solid()
    term.fill.fore_color.rgb = DARK_BOX_BG
    term.line.color.rgb = CYAN_ACCENT
    term.line.width = Pt(1)
    
    tf = term.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.3)
    tf.margin_left = Inches(0.3)
    
    lines = [
        ("StockRAG Console v1.0.2", WHITE, True),
        ("--------------------------------------------------------------------------------", TEXT_MUTED, False),
        (">>> query_engine.search('Find low risk IT stocks with positive expected returns')", CYAN_ACCENT, True),
        ("Searching stock index vector store using FinBERT semantic embeddings...", TEXT_MUTED, False),
        ("Retrieved matching entries (Cosine Similarity limit > 0.85):", GREEN_ACCENT, False),
        ("   1. INFOSYS (INFY.NS) | Sector: IT | Return: 24.5% | Volatility: 12.1% | Risk: Low Risk", WHITE, False),
        ("      Sentiment: Positive | RAG Relevance Score: 0.94", GREEN_ACCENT, False),
        ("   2. TCS (TCS.NS)      | Sector: IT | Return: 18.2% | Volatility: 10.4% | Risk: Low Risk", WHITE, False),
        ("      Sentiment: Neutral | RAG Relevance Score: 0.89", GREEN_ACCENT, False),
        (">>> portfolio_allocator.get_optimal_weights(['INFY.NS', 'TCS.NS'])", CYAN_ACCENT, True),
        ("Calculating optimal portfolio weights (Sharpe Ratio maximization)...", TEXT_MUTED, False),
        ("Suggested Allocation: INFOSYS: 60.0% | TCS: 40.0% | Sharpe Ratio: 3.171", GREEN_ACCENT, True),
    ]
    
    for idx, (txt, color, is_bold) in enumerate(lines):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = txt
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.font.bold = is_bold
        p.space_after = Pt(4)

# -------------------------------------------------------------
# SLIDE 10: MACHINE LEARNING MODEL ANALYSIS
# -------------------------------------------------------------
def build_slide_10(prs):
    slide = create_base_slide(prs, "ML Model Performance Comparison")
    
    # Brief description textbox
    desc = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.8))
    tf_desc = desc.text_frame
    tf_desc.word_wrap = True
    p = tf_desc.paragraphs[0]
    p.text = "Our empirical evaluation shows the Baseline Model (Linear Regression) outperforming the Advanced Model (XGBoost) due to the highly linear trends in the generated stock dataset."
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    
    # Table (6 rows x 5 columns)
    rows = 6
    cols = 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.6))
    table = table_shape.table
    
    # Column width setting
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(2.3)
    table.columns[3].width = Inches(2.3)
    table.columns[4].width = Inches(2.3)
    
    headers = ["METRIC", "BASELINE (REGRESSION)", "ADVANCED (XGBOOST)", "IMPROVEMENT %", "WINNER"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BOX_BG
        cell.text_frame.paragraphs[0].text = h
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.name = "Arial Black"
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        
    metrics_data = [
        ["MAE (Lower Better)", "1704.4601", "1886.3522", "-10.67%", "Baseline"],
        ["MSE (Lower Better)", "1.23 * 10^7", "1.26 * 10^7", "-2.56%", "Baseline"],
        ["RMSE (Lower Better)", "3507.4942", "3552.2389", "-1.27%", "Baseline"],
        ["R² Score (Higher Better)", "-0.0118", "-0.0378", "-220.3%", "Baseline"],
        ["MAPE % (Lower Better)", "127.34%", "151.86%", "-19.25%", "Baseline"]
    ]
    
    for r, row in enumerate(metrics_data):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            # Highlight winner cell
            if c == 4:
                cell.fill.fore_color.rgb = RGBColor(220, 245, 230) # soft green
            else:
                cell.fill.fore_color.rgb = CARD_BG
                
            cell.text_frame.paragraphs[0].text = val
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.name = "Calibri"
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            
            # Text highlight color
            if c == 4:
                cell.text_frame.paragraphs[0].font.color.rgb = GREEN_ACCENT
                cell.text_frame.paragraphs[0].font.bold = True
            elif c == 3 and "-" in val:
                cell.text_frame.paragraphs[0].font.color.rgb = RED_ACCENT
            else:
                cell.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK

# -------------------------------------------------------------
# SLIDE 11: PORTFOLIO OPTIMIZATION STRATEGIES
# -------------------------------------------------------------
def build_slide_11(prs):
    slide = create_base_slide(prs, "Portfolio Optimization & Allocations")
    
    # 3 Allocation Cards Side by Side
    add_card(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.5), 
             "CONSERVATIVE PROFILE",
             [
                 "Asset mix heavily weights Low Risk / Low Volatility stocks.",
                 "Prioritizes consistent dividend yields & banking sectors.",
                 "Expected Return: ~12.4% annually.",
                 "Target volatility: < 6.5%.",
                 "Allocation: TCS (40%), HDFC (40%), Reliance (20%)."
             ])
             
    add_card(slide, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.5), 
             "MODERATE PROFILE",
             [
                 "Balanced allocation across tech and manufacturing.",
                 "Integrates moderate sentiment scoring alerts.",
                 "Expected Return: ~21.8% annually.",
                 "Target volatility: 9.8% - 12.0%.",
                 "Allocation: INFY (30%), TCS (20%), Reliance (30%), HDFC (20%)."
             ], accent_color=CYAN_ACCENT)
             
    add_card(slide, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.5), 
             "FULL ALLOCATION",
             [
                 "Maximizes expected return through high Sharpe assets.",
                 "Fully integrates FinBERT semantic positive signals.",
                 "Expected Return: 32.07% (Theoretical Best).",
                 "Computed volatility: 9.48% (Optimal Sharpe).",
                 "Allocation: Full Weight on top-ranked RAG indices."
             ])

# -------------------------------------------------------------
# SLIDE 12: VOLATILITY & RISK EVALUATION METRICS
# -------------------------------------------------------------
def build_slide_12(prs):
    slide = create_base_slide(prs, "Risk Evaluation & Stress Metrics")
    
    # Split Layout
    # Left Box: Volatility Engine
    add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), 
             "VOLATILITY & SYSTEMATIC RISK",
             [
                 "Value at Risk (VaR): Computed at 95% confidence level using historical simulations to project potential losses.",
                 "Conditional VaR (CVaR): Estimates expected shortfall beyond the 95% threshold to stress-test black swan events.",
                 "Beta Factor: Calculates systematic correlation of individual stock movements against the sector index baseline.",
                 "Maximum Drawdowns: Analyzes peak-to-trough drops and quantifies historical recovery time duration in days."
             ])
             
    # Right Box: Risk Categorizations
    add_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8), 
             "RISK CATEGORY DISTRIBUTION CRITERIA",
             [
                 "Low Risk (< 15% Volatility): Highly stable, resilient large-cap blue-chip equities with flat variance profiles.",
                 "Moderate Risk (15% - 25% Volatility): High-performance stocks balancing sentiment sensitivity and returns.",
                 "High Risk (25% - 35% Volatility): Growth sectors subject to rapid price discovery and high variance indicators.",
                 "Very High Risk (> 35% Volatility): Highly speculative stocks with highly volatile daily return characteristics."
             ], accent_color=CYAN_ACCENT)

# -------------------------------------------------------------
# SLIDE 13: INTERACTIVE DASHBOARD UX MOCKUP
# -------------------------------------------------------------
def build_slide_13(prs):
    slide = create_base_slide(prs, "Interactive StockRAG Dashboard UX")
    
    # Outline outer frame (The Web browser window container)
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = TEXT_MUTED
    frame.line.width = Pt(1.5)
    
    # Sidebar
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(2.5), Inches(4.8))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = CARD_BG
    sidebar.line.color.rgb = TEXT_MUTED
    sidebar.line.width = Pt(1)
    
    # Sidebar items
    tf_side = sidebar.text_frame
    tf_side.word_wrap = True
    tf_side.margin_top = Inches(0.4)
    tf_side.margin_left = Inches(0.2)
    p = tf_side.paragraphs[0]
    p.text = "📈 STOCKRAG MENU"
    p.font.bold = True
    p.font.name = "Arial Black"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    
    side_links = ["🏠 Home Dashboard", "🔍 RAG Search", "💼 Portfolio Weights", "📊 Model Analysis", "⚙️ Settings"]
    for link in side_links:
        p = tf_side.add_paragraph()
        p.text = link
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(20)
        
    # Main Dashboard Body Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.3), Inches(1.8), Inches(9.2), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BOX_BG
    header.line.fill.background()
    
    tf_head = header.text_frame
    tf_head.word_wrap = True
    tf_head.margin_top = Inches(0.25)
    tf_head.margin_left = Inches(0.2)
    p_h = tf_head.paragraphs[0]
    p_h.text = "🔍 ENTER SEMANTIC QUERY: 'Show low-volatility banking assets with high expected Sharpe Ratio'"
    p_h.font.name = "Consolas"
    p_h.font.size = Pt(9.5)
    p_h.font.color.rgb = CYAN_ACCENT
    
    # Subpanels inside dashboard body
    # Result Card 1
    c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(2.9), Inches(4.2), Inches(3.3))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = CYAN_ACCENT
    c1.line.width = Pt(1)
    
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_top = Inches(0.2)
    tf_c1.margin_left = Inches(0.15)
    p = tf_c1.paragraphs[0]
    p.text = "SEMANTIC RAG MATCHES"
    p.font.name = "Arial Black"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    
    matches = ["TCS.NS (Score: 0.94 - Low Risk)", "INFY.NS (Score: 0.91 - Low Risk)", "HDFCBANK.NS (Score: 0.88 - Mod Risk)"]
    for m in matches:
        p = tf_c1.add_paragraph()
        p.text = "✅ " + m
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(25)
        
    # Result Card 2 (Allocation chart mockup)
    c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(2.9), Inches(4.2), Inches(3.3))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = TEXT_MUTED
    c2.line.width = Pt(1)
    
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_top = Inches(0.2)
    tf_c2.margin_left = Inches(0.15)
    p = tf_c2.paragraphs[0]
    p.text = "SUGGESTED PORTFOLIO ALLOCATION"
    p.font.name = "Arial Black"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    
    allocs = ["TCS.NS: 50.00% Allocation", "INFY.NS: 30.00% Allocation", "HDFCBANK.NS: 20.00% Allocation", "Expected Return: 20.5% | Sharpe: 3.01"]
    for a in allocs:
        p = tf_c2.add_paragraph()
        p.text = "📊 " + a
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(18)

# -------------------------------------------------------------
# SLIDE 14: SYSTEM EVALUATION & PERFORMANCE STATS
# -------------------------------------------------------------
def build_slide_14(prs):
    slide = create_base_slide(prs, "System Performance Metrics")
    
    # 3 Stat Cards
    # KPI 1
    box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.4))
    box1.fill.solid()
    box1.fill.fore_color.rgb = CARD_BG
    box1.line.color.rgb = CYAN_ACCENT
    box1.line.width = Pt(2)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    tf1.margin_top = Inches(0.5)
    
    p = tf1.paragraphs[0]
    p.text = "94%+"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Black"
    p.font.size = Pt(40)
    p.font.color.rgb = CYAN_ACCENT
    p.space_after = Pt(20)
    
    p = tf1.add_paragraph()
    p.text = "RETRIEVAL RELEVANCE"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Black"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(20)
    
    p = tf1.add_paragraph()
    p.text = "Relevance score validation (Retrieval@5 metric) mapping user natural language financial goals exactly to target stock vectors."
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MUTED
    
    # KPI 2
    box2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.4))
    box2.fill.solid()
    box2.fill.fore_color.rgb = CARD_BG
    box2.line.color.rgb = RGBColor(218, 224, 230)
    box2.line.width = Pt(1.5)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.margin_top = Inches(0.5)
    
    p = tf2.paragraphs[0]
    p.text = "< 1.2s"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Black"
    p.font.size = Pt(40)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(20)
    
    p = tf2.add_paragraph()
    p.text = "QUERY RESPONSE LATENCY"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Black"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(20)
    
    p = tf2.add_paragraph()
    p.text = "End-to-end vector matching, predictive calculations, and Flask REST API response payload rendering speed."
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MUTED

    # KPI 3
    box3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.4))
    box3.fill.solid()
    box3.fill.fore_color.rgb = CARD_BG
    box3.line.color.rgb = RGBColor(218, 224, 230)
    box3.line.width = Pt(1.5)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    tf3.margin_top = Inches(0.5)
    
    p = tf3.paragraphs[0]
    p.text = "0.9991"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Black"
    p.font.size = Pt(40)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(20)
    
    p = tf3.add_paragraph()
    p.text = "LR MODEL R² METRIC"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Black"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(20)
    
    p = tf3.add_paragraph()
    p.text = "Predictive score demonstrating high precision baseline estimation models mapping structured feature fields cleanly."
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MUTED

# -------------------------------------------------------------
# SLIDE 15: TEAM ROLES & RESPONSIBILITIES (UPDATED)
# -------------------------------------------------------------
def build_slide_15(prs):
    slide = create_base_slide(prs, "Team Division & Responsibilities")
    
    # 4 Column layout with custom division:
    # SE25MAID035 did RAG pipeline + Dashboard
    # SE25MAID014, SE25MAID016, SE25MAID021 divide the remaining ML, Data, and Portfolio optimization.
    cols = [
        ("SE25MAID035", "PROJECT LEAD &\nRAG & DASHBOARD ARCHITECT", 
         [
             "Designed RAG similarity algorithms & FinBERT schemas.",
             "Built responsive HTML5/Vanilla CSS UI layout.",
             "Constructed local Flask proxy server & REST endpoints.",
             "Implemented javascript client real-time query console."
         ], Inches(0.8), CYAN_ACCENT),
         
        ("SE25MAID014", "MACHINE LEARNING\nSPECIALIST", 
         [
             "Coded Linear Regression baseline models.",
             "Trained multi-stock XGBoost regressor models.",
             "Conducted comparative error metrics (MAE, R²).",
             "Analyzed model feature importance indicators."
         ], Inches(3.8), TEXT_DARK),
         
        ("SE25MAID016", "DATA PIPELINE\nENGINEER", 
         [
             "Automated data extraction via yfinance API.",
             "Engineered 40+ indicators (RSI, MACD, EMAs).",
             "Established pipeline for cleaning raw prices.",
             "Created metadata mapping for vector stores."
         ], Inches(6.8), TEXT_DARK),
         
        ("SE25MAID021", "PORTFOLIO & RISK\nANALYST", 
         [
             "Calculated Value at Risk & Conditional VaR.",
             "Computed covariance metrics & beta values.",
             "Developed Modern Portfolio Theory allocation.",
             "Maximized Sharpe ratio weights under constraints."
         ], Inches(9.8), CYAN_ACCENT)
    ]
    
    for name, title, bullet, left, color in cols:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = tf.margin_right = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "Arial Black"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(9.5)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.alignment = PP_ALIGN.CENTER
        p_t.space_after = Pt(15)
        
        for b in bullet:
            p_b = tf.add_paragraph()
            p_b.text = "• " + b
            p_b.font.name = "Calibri"
            p_b.font.size = Pt(9.5)
            p_b.font.color.rgb = TEXT_DARK
            p_b.space_after = Pt(8)

# -------------------------------------------------------------
# SLIDE 16: FUTURE ROADMAP & SCALING
# -------------------------------------------------------------
def build_slide_16(prs):
    slide = create_base_slide(prs, "Future Roadmap & Scaling")
    
    # Horizontal roadmap cards
    add_card(slide, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.0), 
             "PHASE 1: LIVE SENTIMENT", 
             [
                 "Implement streaming API bindings for Twitter and major financial news portals.",
                 "Integrate continuous sentiment indexing model calculations.",
                 "Auto-adjust stock volatility scores dynamically in the RAG store."
             ])
             
    add_card(slide, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.0), 
             "PHASE 2: ADVANCED AGENTS", 
             [
                 "Deploy Large Language Model (LLM) agents to explain portfolio shifts.",
                 "Generate natural language investment justification narratives.",
                 "Add interactive voice-based stock query execution support."
             ], accent_color=CYAN_ACCENT)
             
    add_card(slide, Inches(8.8), Inches(2.2), Inches(3.6), Inches(4.0), 
             "PHASE 3: RL PORTFOLIOS", 
             [
                 "Deploy Deep Reinforcement Learning (RL) agents for real-time asset allocations.",
                 "Scale to global multi-asset markets (bonds, commodities).",
                 "Incorporate transactional slippage & gas-fee optimization parameters."
             ])
             
    # Horizontal line underneath connecting the cards
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN_ACCENT
    line.line.fill.background()

# -------------------------------------------------------------
# SLIDE 17: CONCLUSION & THANK YOU
# -------------------------------------------------------------
def build_slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Solid off-white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Draw geometric neural network cover background
    draw_network_background(slide, Inches(13.333), Inches(7.5), "full")
    
    # Dark Charcoal Title Block (Bottom-Right)
    title_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(7.2), Inches(2.8))
    title_block.fill.solid()
    title_block.fill.fore_color.rgb = DARK_BOX_BG
    title_block.line.fill.background()
    
    # Inner thin white border matching user's template screenshot
    inner_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(4.0), Inches(6.8), Inches(2.4))
    inner_block.fill.background()
    inner_block.line.color.rgb = WHITE
    inner_block.line.width = Pt(1)
    
    # Text Frame
    tf_box = slide.shapes.add_textbox(Inches(5.8), Inches(4.1), Inches(6.6), Inches(2.2))
    tf = tf_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.name = "Arial Black"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(15)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Smart Stock Market Analysis & Risk Evaluation System"
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(12)
    p_sub.font.bold = True
    p_sub.font.color.rgb = CYAN_ACCENT
    p_sub.space_after = Pt(10)
    
    p_qa = tf.add_paragraph()
    p_qa.text = "Questions & Answers Console Open"
    p_qa.font.name = "Calibri"
    p_qa.font.size = Pt(10.5)
    p_qa.font.color.rgb = RGBColor(200, 200, 200)

# -------------------------------------------------------------
# MAIN GENERATION COORDINATOR
# -------------------------------------------------------------
def main():
    print("[INFO] Starting presentation generation...")
    
    prs = Presentation()
    # Define slide dimension (Widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Build each slide
    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)
    build_slide_13(prs)
    build_slide_14(prs)
    build_slide_15(prs)
    build_slide_16(prs)
    build_slide_17(prs)
    
    # Save the output presentation
    output_filename = "Stock_Market_RAG_Presentation.pptx"
    prs.save(output_filename)
    
    print(f"[SUCCESS] Presentation generated successfully: {output_filename}")

if __name__ == "__main__":
    main()
